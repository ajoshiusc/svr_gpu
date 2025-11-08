#!/usr/bin/env python
"""
Generate a simple PowerPoint (.pptx) from SMS_VALIDATION_REPORT.md.

Requirements:
  pip install python-pptx markdown

Output:
  SMS_VALIDATION_REPORT.pptx (in project root)
"""
from pathlib import Path
import re
import sys

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_PARAGRAPH_ALIGNMENT


ROOT = Path(__file__).resolve().parents[1]
REPORT_MD = ROOT / "SMS_VALIDATION_REPORT.md"
OUTPUT_PPTX = ROOT / "SMS_VALIDATION_REPORT.pptx"


def parse_markdown_headings(md_text: str):
    """Parse markdown into a simple list of (level, title, content) sections.

    We keep only top (h1) and section (h2) headings as slides.
    Content under each section is truncated/cleaned to bullet points.
    """
    lines = md_text.splitlines()
    sections = []
    current = {"level": 1, "title": "", "content": []}

    def push_current():
        if current["title"] or current["content"]:
            sections.append({
                "level": current["level"],
                "title": current["title"].strip(),
                "content": "\n".join(current["content"]).strip()
            })

    for ln in lines:
        m = re.match(r"^(#{1,6})\s+(.*)$", ln)
        if m:
            # new heading
            level = len(m.group(1))
            title = m.group(2).strip()
            if level <= 2:
                push_current()
                current = {"level": level, "title": title, "content": []}
            # skip deeper headings for slide titles; keep as content bullets
            else:
                current["content"].append(f"- {title}")
        else:
            current["content"].append(ln)

    push_current()
    return sections


def bullets_from_text(text: str):
    """Convert a block of text into a list of bullet lines (trimmed, non-empty)."""
    bullets = []
    for raw in text.splitlines():
        t = raw.strip()
        if not t:
            continue
        # Remove markdown fences and code markers for brevity
        if t.startswith("```"):
            continue
        # Collapse long lines
        t = re.sub(r"\s+", " ", t)
        # Strip leading markdown bullets
        t = re.sub(r"^[-*]\\s?", "", t)
        # Limit bullet length
        if len(t) > 180:
            t = t[:177] + "..."
        bullets.append(t)
        # Keep slides concise
        if len(bullets) >= 10:
            break
    return bullets


def add_title_slide(prs: Presentation, title: str, subtitle: str = ""):
    layout = prs.slide_layouts[0]  # Title slide
    slide = prs.slides.add_slide(layout)
    slide.shapes.title.text = title
    if subtitle:
        slide.placeholders[1].text = subtitle


def add_bullet_slide(prs: Presentation, title: str, bullets: list[str]):
    layout = prs.slide_layouts[1]  # Title and Content
    slide = prs.slides.add_slide(layout)
    slide.shapes.title.text = title
    tf = slide.placeholders[1].text_frame
    tf.clear()
    for i, b in enumerate(bullets):
        if i == 0:
            p = tf.paragraphs[0]
            p.text = b
        else:
            p = tf.add_paragraph()
            p.text = b
        p.level = 0
        p.font.size = Pt(18)
        p.alignment = PP_PARAGRAPH_ALIGNMENT.LEFT


def main():
    if not REPORT_MD.exists():
        print(f"Report not found: {REPORT_MD}")
        sys.exit(1)

    md = REPORT_MD.read_text(encoding="utf-8")
    sections = parse_markdown_headings(md)

    prs = Presentation()

    # Title slide
    title = "SMS-Aware SVR: Validation Summary"
    subtitle = "Generated from SMS_VALIDATION_REPORT.md"
    add_title_slide(prs, title, subtitle)

    # Build slides from H2 sections primarily; include H1 Executive Summary if present
    for sec in sections:
        lvl = sec["level"]
        title = sec["title"] or "Summary"
        content = sec["content"]

        # Skip the top H1 title (already used for overall title slide)
        if lvl == 1 and "Executive Summary" not in title:
            continue

        # Build bullets
        bullets = bullets_from_text(content)
        if not bullets:
            continue
        add_bullet_slide(prs, title, bullets)

    prs.save(OUTPUT_PPTX)
    print(f"Saved: {OUTPUT_PPTX}")


if __name__ == "__main__":
    main()
